#!/usr/bin/env python3
"""
PPT Generator Skill
生成精美的PowerPoint演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import sys
import os

def create_title_slide(prs, title, subtitle=""):
    """创建标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 20, 25)  # 深色背景
    shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(154, 160, 166)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, bullet_points, layout="default"):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(22, 119, 255)
    bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.space_after = Pt(12)
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """创建双栏布局页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    
    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(10)
    
    # 右栏
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(10)
    
    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.2), Inches(0.02), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.fill.background()
    
    return slide

def generate_ppt(data):
    """生成PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    title = data.get('title', '演示文稿')
    subtitle = data.get('subtitle', '')
    create_title_slide(prs, title, subtitle)
    
    # 内容页
    slides = data.get('slides', [])
    for slide_data in slides:
        slide_type = slide_data.get('type', 'content')
        slide_title = slide_data.get('title', '')
        
        if slide_type == 'content':
            points = slide_data.get('points', [])
            create_content_slide(prs, slide_title, points)
        elif slide_type == 'two-column':
            left = slide_data.get('left', [])
            right = slide_data.get('right', [])
            create_two_column_slide(prs, slide_title, left, right)
    
    # 保存
    output_file = data.get('output', f"presentation_{os.getpid()}.pptx")
    prs.save(output_file)
    
    return output_file

def main():
    """主函数"""
    if len(sys.argv) < 2:
        print(json.dumps({'error': 'Usage: python ppt_generator.py \'<json_data>\''}))
        sys.exit(1)
    
    try:
        data = json.loads(sys.argv[1])
        output_file = generate_ppt(data)
        
        result = {
            'success': True,
            'output': output_file,
            'message': f'PPT generated: {output_file}'
        }
        print(json.dumps(result, ensure_ascii=False))
        
    except Exception as e:
        result = {
            'success': False,
            'error': str(e)
        }
        print(json.dumps(result, ensure_ascii=False))
        sys.exit(1)

if __name__ == '__main__':
    main()
